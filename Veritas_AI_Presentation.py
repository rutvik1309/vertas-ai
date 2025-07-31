#!/usr/bin/env python3
"""
PowerPoint Presentation Generator for Veritas AI Project
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
import os

def create_veritas_ai_presentation():
    """Create a comprehensive PowerPoint presentation for Veritas AI"""
    
    # Create presentation
    prs = Presentation()
    
    # Define colors
    primary_color = RGBColor(0, 123, 255)  # Blue
    accent_color = RGBColor(255, 193, 7)   # Yellow
    dark_color = RGBColor(52, 58, 64)      # Dark gray
    success_color = RGBColor(40, 167, 69)  # Green
    warning_color = RGBColor(255, 193, 7)  # Yellow
    danger_color = RGBColor(220, 53, 69)   # Red
    
    # Slide 1: Title Slide
    slide_layout = prs.slide_layouts[0]  # Title slide
    slide = prs.slides.add_slide(slide_layout)
    
    title = slide.shapes.title
    subtitle = slide.placeholders[1]
    
    title.text = "VERITAS AI"
    subtitle.text = "News Fact Checker\nAI-Powered Misinformation Detection System"
    
    # Style title
    title.text_frame.paragraphs[0].font.size = Pt(44)
    title.text_frame.paragraphs[0].font.color.rgb = primary_color
    title.text_frame.paragraphs[0].font.bold = True
    
    # Style subtitle
    subtitle.text_frame.paragraphs[0].font.size = Pt(20)
    subtitle.text_frame.paragraphs[0].font.color.rgb = dark_color
    
    # Slide 2: Agenda
    slide_layout = prs.slide_layouts[1]  # Title and content
    slide = prs.slides.add_slide(slide_layout)
    
    title = slide.shapes.title
    content = slide.placeholders[1]
    
    title.text = "Presentation Agenda"
    title.text_frame.paragraphs[0].font.color.rgb = primary_color
    
    content_text = """• Project Objective
• Scope Definition
• Target Audience
• Business Case
• Technologies Used
• Challenges & Solutions
• System Thresholds
• Key Achievements"""
    
    content.text = content_text
    
    # Slide 3: Objective
    slide_layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(slide_layout)
    
    title = slide.shapes.title
    content = slide.placeholders[1]
    
    title.text = "Project Objective"
    title.text_frame.paragraphs[0].font.color.rgb = primary_color
    
    content_text = """🎯 Primary Goal:
Develop an AI-powered news fact-checking system that automatically analyzes and verifies the authenticity of news articles, videos, and multimedia content using machine learning and natural language processing.

Key Objectives:
• Automated Fact-Checking: Real-time analysis of news content
• Multi-Format Support: Text, YouTube videos, images, audio, documents
• Credibility Assessment: Evaluate source reliability and content quality
• User-Friendly Interface: Chrome extension for seamless integration
• Scalable Architecture: Handle multiple users with API key rotation"""
    
    content.text = content_text
    
    # Slide 4: In Scope
    slide_layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(slide_layout)
    
    title = slide.shapes.title
    content = slide.placeholders[1]
    
    title.text = "In Scope"
    title.text_frame.paragraphs[0].font.color.rgb = success_color
    
    content_text = """✅ Core Features:
• News Classification: ML model with confidence scores
• Multi-Media Analysis: Text, images, audio, video, documents
• YouTube Integration: Specialized analysis with transcript extraction
• Source Credibility: Reputable source detection
• API Key Management: Automatic rotation across 8 keys
• Chrome Extension: Browser-based interface
• Web Application: Full-featured with authentication
• Conversation Memory: Learning from interactions
• Performance Monitoring: Real-time tracking

✅ Technical Components:
• Backend: Flask web application with RESTful APIs
• Frontend: Chrome extension with modern UI
• Database: SQLite with user management
• ML Pipeline: Scikit-learn models
• AI Integration: Google Gemini API
• Deployment: Docker containerization"""
    
    content.text = content_text
    
    # Slide 5: Out of Scope
    slide_layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(slide_layout)
    
    title = slide.shapes.title
    content = slide.placeholders[1]
    
    title.text = "Out of Scope"
    title.text_frame.paragraphs[0].font.color.rgb = danger_color
    
    content_text = """❌ Not Included:
• Real-time Video Streaming Analysis
• Social Media Integration
• Multi-language Support (English only)
• Mobile Application
• Advanced Video Processing
• Real-time News Monitoring
• Legal Compliance
• Advanced NLP"""
    
    content.text = content_text
    
    # Slide 6: Target Audience
    slide_layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(slide_layout)
    
    title = slide.shapes.title
    content = slide.placeholders[1]
    
    title.text = "Target Audience"
    title.text_frame.paragraphs[0].font.color.rgb = primary_color
    
    content_text = """👥 Primary Users:
• General Public: Individuals seeking to verify news authenticity
• Students: Academic research and fact-checking
• Journalists: Quick verification of sources and claims
• Researchers: Academic fact-checking and validation
• Content Creators: Verification before sharing

👥 Secondary Users:
• Educational Institutions: Teaching media literacy
• News Organizations: Internal fact-checking tools
• Government Agencies: Public information verification
• Business Professionals: Corporate communication verification"""
    
    content.text = content_text
    
    # Slide 7: Business Case
    slide_layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(slide_layout)
    
    title = slide.shapes.title
    content = slide.placeholders[1]
    
    title.text = "Business Case"
    title.text_frame.paragraphs[0].font.color.rgb = primary_color
    
    content_text = """📊 Market Need:
• Misinformation Crisis: 64% of adults encounter fake news regularly
• Digital Literacy Gap: Need for accessible fact-checking tools
• Information Overload: Users struggle to verify content authenticity
• Media Literacy: Growing demand for critical thinking tools

💰 Value Proposition:
• Accessibility: Easy-to-use Chrome extension
• Accuracy: ML + AI hybrid approach
• Comprehensive: Multi-format support
• Scalable: API rotation system
• Educational: Develops critical thinking skills

💼 Revenue Potential:
• Freemium Model: Basic free, premium for advanced features
• API Services: B2B fact-checking API
• Educational Licensing: School partnerships
• Enterprise Solutions: Custom deployments"""
    
    content.text = content_text
    
    # Slide 8: Technologies Used
    slide_layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(slide_layout)
    
    title = slide.shapes.title
    content = slide.placeholders[1]
    
    title.text = "Technologies Used"
    title.text_frame.paragraphs[0].font.color.rgb = primary_color
    
    content_text = """🔧 Backend Technologies:
• Python 3.11, Flask 3.1.0, SQLAlchemy
• Flask-Login, Flask-CORS

🤖 Machine Learning & AI:
• Scikit-learn 1.7.1, Google Gemini API
• NLTK 3.8.1, Joblib 1.4.2, NumPy 1.24.4

📊 Data Processing:
• Newspaper3k, BeautifulSoup4, Pillow 11.0.0
• PyPDF2, python-docx

🎥 Media Processing:
• yt-dlp, youtube-transcript-api
• SpeechRecognition, pytesseract, moviepy

🌐 Frontend & Extension:
• Chrome Extension API, JavaScript ES6
• HTML5/CSS3, Chrome Storage API

🚀 Deployment & Infrastructure:
• Docker, Gunicorn, Nginx
• ChromaDB, Redis"""
    
    content.text = content_text
    
    # Slide 9: Challenges & Solutions (Part 1)
    slide_layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(slide_layout)
    
    title = slide.shapes.title
    content = slide.placeholders[1]
    
    title.text = "Challenges & Solutions (Part 1)"
    title.text_frame.paragraphs[0].font.color.rgb = warning_color
    
    content_text = """🔴 Challenge 1: API Rate Limiting
Problem: Google Gemini API free tier limited to 50 requests/day
Solution: 
• Automatic API key rotation across 8 keys
• Performance tracking system
• Intelligent load balancing
Result: 400 requests/day capacity

🔴 Challenge 2: YouTube Content Analysis
Problem: ML model incorrectly classified reputable videos as "Fake"
Solution:
• Hybrid analysis system for YouTube content
• Reputable source detection (CBC, BBC, Reuters)
• Natural language processing for metadata
• Sensationalist content detection
Result: Accurate classification of reputable vs. fake content

🔴 Challenge 3: Model Compatibility
Problem: Scikit-learn version conflicts
Solution:
• Model version compatibility scripts
• Fallback mechanisms
• Comprehensive error handling
Result: Stable model loading across environments"""
    
    content.text = content_text
    
    # Slide 10: Challenges & Solutions (Part 2)
    slide_layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(slide_layout)
    
    title = slide.shapes.title
    content = slide.placeholders[1]
    
    title.text = "Challenges & Solutions (Part 2)"
    title.text_frame.paragraphs[0].font.color.rgb = warning_color
    
    content_text = """🔴 Challenge 4: Extension-Backend Communication
Problem: Chrome extension connection errors
Solution:
• Fixed CORS configuration
• Proper conversation state management
• Error handling for file uploads
Result: Seamless extension-backend communication

🔴 Challenge 5: Multi-Format Processing
Problem: Supporting various file formats
Solution:
• Multiple processing libraries integration
• OCR for image text extraction
• Audio-to-text conversion
Result: Comprehensive multi-format support

🔴 Challenge 6: Scalability
Problem: Single API key bottleneck
Solution:
• Automatic API key rotation system
• Performance monitoring and tracking
• Caching mechanisms
Result: Production-ready scalable architecture"""
    
    content.text = content_text
    
    # Slide 11: System Thresholds
    slide_layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(slide_layout)
    
    title = slide.shapes.title
    content = slide.placeholders[1]
    
    title.text = "System Thresholds"
    title.text_frame.paragraphs[0].font.color.rgb = primary_color
    
    content_text = """📊 Performance Thresholds:
• Response Time: < 30 seconds for complex analysis
• Accuracy: > 85% for reputable source detection
• API Success Rate: > 90% for key rotation system
• Uptime: > 99% for production deployment
• Concurrent Users: Support for 100+ simultaneous users

📊 Technical Thresholds:
• File Size Limits: Images (10MB), Videos (100MB), Documents (50MB)
• Text Length: 10,000 characters max per analysis
• API Requests: 60 requests/minute per key
• Daily Capacity: 400 requests/day (8 keys × 50 requests)

📊 Quality Thresholds:
• Confidence Score: > 70% for reliable predictions
• Source Verification: 100% for known reputable sources
• Error Recovery: Automatic fallback for failed API calls
• Data Privacy: No personal data storage in analysis

📊 User Experience Thresholds:
• Extension Load Time: < 3 seconds
• Analysis Completion: < 30 seconds for most content
• UI Responsiveness: < 1 second for user interactions
• Error Handling: Graceful degradation with helpful messages"""
    
    content.text = content_text
    
    # Slide 12: Key Achievements
    slide_layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(slide_layout)
    
    title = slide.shapes.title
    content = slide.placeholders[1]
    
    title.text = "Key Achievements"
    title.text_frame.paragraphs[0].font.color.rgb = success_color
    
    content_text = """🏆 Major Accomplishments:

✅ Multi-Format Support
• Text, images, audio, video, documents
• Comprehensive content analysis capabilities

✅ Hybrid AI/ML System
• Combines ML classification with AI reasoning
• Enhanced accuracy and reliability

✅ Scalable Architecture
• API rotation for high-volume usage
• Production-ready infrastructure

✅ User-Friendly Interface
• Chrome extension for easy access
• Intuitive web application

✅ Educational Value
• Helps develop critical thinking skills
• Promotes media literacy

✅ Comprehensive Analysis
• Source credibility assessment
• Content quality evaluation
• Real-time fact-checking capabilities

🎯 Result: A sophisticated, production-ready fact-checking system with innovative solutions to real-world challenges in the fight against misinformation."""
    
    content.text = content_text
    
    # Slide 13: Thank You
    slide_layout = prs.slide_layouts[0]  # Title slide
    slide = prs.slides.add_slide(slide_layout)
    
    title = slide.shapes.title
    subtitle = slide.placeholders[1]
    
    title.text = "Thank You!"
    subtitle.text = "Questions & Discussion\n\nVeritas AI - Fighting Misinformation with AI"
    
    # Style title
    title.text_frame.paragraphs[0].font.size = Pt(44)
    title.text_frame.paragraphs[0].font.color.rgb = success_color
    title.text_frame.paragraphs[0].font.bold = True
    
    # Style subtitle
    subtitle.text_frame.paragraphs[0].font.size = Pt(20)
    subtitle.text_frame.paragraphs[0].font.color.rgb = dark_color
    
    # Save the presentation
    filename = "Veritas_AI_Presentation.pptx"
    prs.save(filename)
    print(f"✅ Presentation created successfully: {filename}")
    print(f"📁 File saved in: {os.getcwd()}")
    
    return filename

if __name__ == "__main__":
    try:
        # Install required package if not available
        import subprocess
        import sys
        
        try:
            from pptx import Presentation
        except ImportError:
            print("Installing python-pptx...")
            subprocess.check_call([sys.executable, "-m", "pip", "install", "python-pptx"])
            from pptx import Presentation
        
        filename = create_veritas_ai_presentation()
        print(f"\n🎉 Your Veritas AI presentation is ready!")
        print(f"📊 Total slides: 13")
        print(f"📋 Content includes: Objective, Scope, Audience, Business Case, Technologies, Challenges, Solutions, Thresholds, and Achievements")
        
    except Exception as e:
        print(f"❌ Error creating presentation: {e}")
        print("💡 Make sure you have write permissions in the current directory") 